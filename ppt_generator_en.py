from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import sys
from pathlib import Path

DOCUMENTS = Path.home() / "Documents"
OUTPUT_FOLDER = DOCUMENTS / "Quality Reports"

OUTPUT_FOLDER.mkdir(
    parents=True,
    exist_ok=True
)


def resource_path(relative_path):

    try:
        base_path = sys._MEIPASS
    except AttributeError:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)


TEMPLATE = resource_path("templates_ppt/Quality_Gate_Template.pptx")


# ==========================================================
# Substitui os marcadores em uma TextFrame
# ==========================================================

def replace_in_text_frame(text_frame, values):

    for paragraph in text_frame.paragraphs:

        for run in paragraph.runs:

            text = run.text

            for key, value in values.items():

                marker = "{{" + key + "}}"

                if marker in text:
                    text = text.replace(marker, str(value))

            run.text = text


# ==========================================================
# Percorre todos os objetos do slide
# ==========================================================

def replace_tags(shape, values):

    if hasattr(shape, "text_frame") and shape.has_text_frame:
        replace_in_text_frame(shape.text_frame, values)

    if shape.has_table:

        table = shape.table

        for row in table.rows:
            for cell in row.cells:
                replace_in_text_frame(cell.text_frame, values)

    if shape.shape_type == 6:

        for shp in shape.shapes:
            replace_tags(shp, values)


# ==========================================================
# Remove todas as imagens do slide
# ==========================================================

def remove_images(slide):

    for shape in list(slide.shapes):

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

            shape._element.getparent().remove(shape._element)


# ==========================================================
# Geração do PPT EM INGLÊS
# ==========================================================

def generate_ppt_en(
    model,
    line,
    process,
    date,
    qty_defect,
    category,
    defect,
    cause,
    action,
    pic,
    image_path=""
):

    prs = Presentation(TEMPLATE)

    values = {

        "MODEL": model,
        "LINE": line,
        "PROCESS": process,
        "DATE": date,
        "QTY_DEFECT": qty_defect,
        "CATEGORY": category,
        "DEFECT": defect,
        "CAUSE": cause,
        "ACTION": action,
        "PIC": pic

    }

    for slide in prs.slides:

        for shape in slide.shapes:
            replace_tags(shape, values)

        remove_images(slide)

        if image_path and os.path.exists(image_path):

            slide.shapes.add_picture(
                image_path,
                left=7.8 * 914400,
                top=2.0 * 914400,
                width=2.8 * 914400,
                height=4.2 * 914400
            )

    filename = OUTPUT_FOLDER / f"{model}_Quality_Report_EN.pptx"

    prs.save(str(filename))

    return str(filename)