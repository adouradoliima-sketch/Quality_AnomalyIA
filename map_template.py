from pptx import Presentation
import os

arquivo = "templates_ppt/Quality_Gate_Template.pptx"

print("=" * 60)
print("QUALITY ANOMALY - TEMPLATE READER")
print("=" * 60)

print("Arquivo encontrado:", os.path.exists(arquivo))

prs = Presentation(arquivo)

print("Quantidade de Slides:", len(prs.slides))

for s, slide in enumerate(prs.slides):

    print("\n")
    print("=" * 60)
    print(f"SLIDE {s}")
    print("=" * 60)

    print("Quantidade de Shapes:", len(slide.shapes))

    for i, shape in enumerate(slide.shapes):

        print("\n----------------------------------------")
        print(f"Shape {i}")
        print("----------------------------------------")
        print("Tipo:", shape.shape_type)

        if hasattr(shape, "text"):

            texto = shape.text.strip()

            if texto != "":
                print("Texto:")
                print(texto)

        if shape.has_table:

            print("\n***** TABELA *****")

            table = shape.table

            for r in range(len(table.rows)):

                for c in range(len(table.columns)):

                    texto = table.cell(r, c).text

                    print(f"Célula [{r},{c}]")

                    print(texto)

                    print("----------------------------")